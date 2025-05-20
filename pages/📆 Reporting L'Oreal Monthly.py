import streamlit as st
import os
import pandas as pd
import numpy as np
import time
from datetime import datetime
import json
import smtplib
import pytz
import requests
from bs4 import BeautifulSoup
from io import BytesIO
from PIL import Image
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from google.cloud import bigquery
from google.oauth2 import service_account
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

st.logo(
    "image/goat_logo.png", size="medium", link=None, icon_image=None
)
st.image("image/goat-agency-logo2.png")
st.header("L'Oreal Monthly Report")
year = st.selectbox(
  'Please select the reporting year',
  ('2025', '2024', '2023')
)
year_ = {
  '2023':23, '2024':24, '2025':25
}

year_map = year_.get(year, "")  # Returns '' if year is not found
year_range = (year_map -1, year_map)

quarter_ = st.multiselect(
  'Please select the reporting quarter',
  ['Quarter 1', 'Quarter 2', 'Quarter 3', 'Quarter 4'], default = ['Quarter 1'], max_selections=2
)
q_map = {
  'Quarter 1': 'Q1', 'Quarter 2': 'Q2', 'Quarter 3': 'Q3', 'Quarter 4': 'Q4'
}
quarter = [q_map[q] for q in quarter_ if q in q_map]

month = st.multiselect(
  'Please select the reporting month',
  ["Jan", "Feb", "Mar", "Apr",
    "May", "Jun", "Jul", "Aug",
    "Sep", "Oct", "Nov", "Dec"], default = ["Jan"]
)

month_map = {
    "Jan": "1", "Feb": "2", "Mar": "3", "Apr": "4",
    "May": "5", "Jun": "6", "Jul": "7", "Aug": "8",
    "Sep": "9", "Oct": "10", "Nov": "11", "Dec": "12"
}

month_num = [month_map[m] for m in month if m in month_map]

division_selection = st.multiselect(
  "Please select the reporting L'Oreal Division",
  ['Select All', 'CPD', 'LDB', 'LLD', 'PPD'], default=['Select All']
)
if 'Select All' in division_selection:
    division = ['CPD', 'LDB', 'LLD', 'PPD']
else: division = division_selection
	
category_selection = st.multiselect(
  "Please select the reporting L'Oreal Category",
  ['Select All', 'Female Skin', 'Hair Care', 'Hair Color', 'Make Up', 'Men Skin', 'Fragrance', 'Pro Hair'], default=['Select All']
)
if 'Select All' in category_selection:
    category = ['Female Skin', 'Hair Care', 'Hair Color', 'Make Up', 'Men Skin', 'Fragrance', 'Pro Hair']
else: category = category_selection

brands = st.multiselect(
    "Please Select max 5 L'Oreal Brands to compare in the report",
	["3CE", "ARM", "CRV", "GAR", "KER", "KIE", "LAN", "LP", "LRP", "MNY", "MXBIO", "OAP", "SHU", "YSL"]
	,default = ["GAR", "LP", "MNY"]
	,max_selections=5
)
# brand full name -> ["BLP Skin", "Garnier", "L'Oreal Paris", "GMN Shampoo Color", "Armani", "Kiehls", "Lancome", "Shu Uemura", "Urban Decay", "YSL", "Cerave", "La Roche Posay", "L'Oreal Professionel", "Matrix", "Biolage", "Kerastase", "Maybelline"]

email_list = st.multiselect(
	"Please select the email address you wish to send this report to",
	("dyah.dinasari@groupm.com", "bandi.wijaya@groupm.com", "alfian.hidayatulloh@groupm.com", "azka.fatharani@groupm.com",
	 "yahya.wahyu@groupm.com", "agusta.jalasena@groupm.com", "ana.maratu@groupm.com", "aldi.firstanto@groupm.com", 
	 "muhammad.ilham@groupm.com")
	)

st.badge(" Optional", icon="⚠️", color="blue")
st.info("You can upload PPT Template for this report, make sure it have 20 slides or more.")
uploaded_file = st.file_uploader("", type=["pptx"])
if uploaded_file is not None:
	st.success("File uploaded successfully!")
else:
	uploaded_file = "pages/Template Deck to Go - L'Oreal Indonesia.pptx"

if st.button("Generate Report", type="primary"):

#--- DATA PROCESSING ---

#CHART FORMATING
	#--------
	def format_number_kmb(n):
		if n is None:
			return ""
		abs_n = abs(n)
		if abs_n >= 1_000_000_000:
			return f"{n / 1_000_000_000:.1f}B"
		elif abs_n >= 1_000_000:
			return f"{n / 1_000_000:.1f}M"
		elif abs_n >= 1_000:
			return f"{n / 1_000:.1f}K"
		else:
			return str(int(n))
	#-------
	def format_title(slide, text, alignment, font_name, font_size, font_bold = False, font_italic = None, font_color = RGBColor(0, 0, 0),left=Pt(75), top=Pt(25), width=Pt(850), height=Pt(70)):
		title_shape = slide.shapes.add_textbox(left=left, top=top, width=width, height=height)
		title_text_frame = title_shape.text_frame
		title_text_frame.text = text
		title_text_frame.word_wrap = True
		for paragraph in title_text_frame.paragraphs:
			paragraph.alignment = alignment
			for run in paragraph.runs:
				run.font.name = font_name
				run.font.bold = font_bold
				run.font.italic = font_italic
				run.font.color.rgb = font_color
				run.font.size = Pt(font_size)
		return title_shape

	def pie_chart(slide,df,x,y,cx,cy,fontsize=9,legend_right = True, chart_title = False, title='',fontsize_title = Pt(20)):
		df.fillna(0, inplace = True) #fill nan
		# Convert the transposed DataFrame into chart data
		chart_data = CategoryChartData()
		# Add the brand names as categories to the chart data
		for i in df.transpose().columns:
			chart_data.add_category(i)
	  	# Add the SOV values as series to the chart data
		for index, row in df.transpose().iterrows():
			chart_data.add_series(index, row.values)
  		
		# create chart  
		chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
	  	
		if chart_title:
			chart.has_title = True
			chart.chart_title.text_frame.text = title
			title_font = chart.chart_title.text_frame.paragraphs[0].font
			title_font.bold = True
			title_font.size = fontsize_title
	    		# chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to 24pt
	    		# chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)  # Set font color to black
		else:
			chart.has_title = False
	
		if chart.has_legend:
			chart.legend.include_in_layout = False
			chart.legend.position = XL_LEGEND_POSITION.RIGHT if legend_right else XL_LEGEND_POSITION.BOTTOM
			chart.legend.font.size = Pt(fontsize)
	
		for series in chart.plots[0].series:
			for i, val in enumerate(series.values):
				if val == 0:
					series.points[i].data_label.has_text_frame = True
				series.data_labels.show_value = True
				series.data_labels.font.size = Pt(fontsize)
				series.data_labels.number_format = '0%'
				series.data_labels.position = XL_LABEL_POSITION.BEST_FIT
				series.data_labels.show_category_name = True
	
		return chart

	def line_marker_chart(slide,df,x,y,cx,cy, legend = True, legend_position = XL_LEGEND_POSITION.RIGHT,
                      data_show = False, chart_title = False, title ="", fontsize = Pt(12),
                      fontsize_title = Pt(14), percentage = False, line_width = Pt(1), smooth=False):
		df.fillna(0, inplace = True) #fill nan
	  # Define chart data
		chart_data = CategoryChartData()

		for i in df.columns:
			chart_data.add_category(i)
		for j, row in df.iterrows():
			chart_data.add_series(j,np.where(row.values == 0, None, row.values))
	
		chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
	
	  # smoothing
		if smooth: # Smooth the lines
			for series in chart.series:
				series.smooth = True

	  # Add legend
		if legend:
			chart.has_legend = True
			chart.legend.include_in_layout = False
			chart.legend.position = legend_position
			chart.legend.font.size = fontsize  # assuming Pt is imported from pptx.util
		else:
			chart.has_legend = False
	
		if data_show:
			for series in chart.plots[0].series:
				for i, val in enumerate(series.values):
					if val == 0:
						series.points[i].data_label.has_text_frame = True
				series.data_labels.show_value = True
				series.data_labels.font.size = fontsize
	      # series.data_labels.number_format = '0.00%'
				series.data_labels.position = XL_LABEL_POSITION.ABOVE
	
		for series in chart.series:
			series.marker.style = XL_MARKER_STYLE.CIRCLE
			series.format.line.width = line_width # custom line width
	
	  # Customize y-axis format
		chart.value_axis.tick_labels.font.size = fontsize  # Set font size for tick labels
	
	  # Set font size for category axis (months)
		chart.category_axis.tick_labels.font.size = fontsize  # Set font size for category axis labels
	
	  # Find the maximum value across all series
		max_value = 0
		for series in chart.plots[0].series:
			try:
				series_max = max(series.values)
			except:
				series_max = 0
			max_value = max(max_value, series_max)
	
		if max_value >= 1000:
			chart.value_axis.tick_labels.number_format = '#,##0'  # add commas to separate thousands
	
		if chart_title:
			chart.chart_title.text_frame.text = title # Set the title text
			title_font = chart.chart_title.text_frame.paragraphs[0].font
			title_font.bold = True
			title_font.size = fontsize_title
		else:
			chart.has_title = False
	
	  # Remove Gridlines (Line Chart Specific)
		value_axis = chart.value_axis
		category_axis = chart.category_axis
		value_axis.has_major_gridlines = False
		value_axis.has_minor_gridlines = False
		category_axis.has_major_gridlines = False
		category_axis.has_minor_gridlines = False
	
		if percentage:
			category_axis.tick_labels.NumberFormat = '0"%"'
	
		return chart

	def table_default(slide, df, left, top, width, height, width_row, height_row, header=True, upper=False, fontsize=10, alignment=PP_ALIGN.LEFT):
		table_data = df.values.tolist()
		
		if header:
			header = df.columns.values.tolist()
			table_data.insert(0,header)
	
		table = slide.shapes.add_table(rows=len(table_data), cols=len(table_data[0]), left=left, top=top, width=width, height=height).table
	
	    # Fill table data
		for i, row in enumerate(table_data):
			for j, val in enumerate(row):
				cell = table.cell(i, j)
				if isinstance(val, (int, float)):  # Check if it's a number
					text = f"{val:,}"  # Format with comma separators
				else:
					text = str(val).upper() if upper else str(val)
				cell.text = text
				for paragraph in cell.text_frame.paragraphs:
					paragraph.alignment = alignment
					for run in paragraph.runs:
						run.font.size = Pt(fontsize)

				cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
			    # Set all margins to 0
				cell.text_frame.margin_left = 0
				cell.text_frame.margin_right = 0
				cell.text_frame.margin_top = 0
				cell.text_frame.margin_bottom = 0
	
	    	# Set the width of each column
			for i, column in enumerate(table.columns):
				column.width = width_row[i]
			# Adjust the height of each row
			for row in table.rows:
				row.height = height_row # Set the height of each row
		
		return table

	def horizontal_bar_chart(slide, df, x, y, cx, cy, legend=True, legend_position=XL_LEGEND_POSITION.RIGHT,
							 data_show=True, chart_title=False, title="", fontsize=Pt(10),
							 fontsize_title=Pt(14), percentage=False, bar_width=Pt(10)):
		df.fillna(0, inplace=True)  # Fill NaN values
		# Define chart data
		chart_data = CategoryChartData()
		for i in df.index:
			chart_data.add_category(i)
		for col in df.columns:
			chart_data.add_series(col, np.where(df[col].values == 0, None, df[col].values))
	
	    # Create horizontal bar chart
		chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
	
	    # Add legend
		if legend:
			chart.has_legend = True
			chart.legend.include_in_layout = False
			chart.legend.position = legend_position
			chart.legend.font.size = fontsize
		else:
			chart.has_legend = False
	
	    # Show data labels
		if data_show:
			for series_idx, series in enumerate(chart.series):
				for point_idx, point in enumerate(series.points):
					val = df.iloc[point_idx, series_idx]
					label_text = format_number_kmb(val) if val != 0 else ""
					point.data_label.text = label_text
					point.data_label.font.size = fontsize
					point.data_label.position = XL_LABEL_POSITION.OUTSIDE_END
	
	
	    # Adjust bar width
		for series in chart.series:
			series.format.line.width = bar_width
	
		# Set chart title
		if chart_title:
			chart.chart_title.text_frame.text = title
			title_font = chart.chart_title.text_frame.paragraphs[0].font
			title_font.bold = True
			title_font.size = fontsize_title
		else:
			chart.has_title = False
	
		# Remove gridlines
		chart.value_axis.has_major_gridlines = False
		chart.value_axis.has_minor_gridlines = False
		chart.category_axis.has_major_gridlines = False
		chart.category_axis.has_minor_gridlines = False

		# Hide X (value) and Y (category) axes
		chart.value_axis.has_major_tick_marks = False
		chart.value_axis.tick_labels.font.size = fontsize
		chart.value_axis.format.line.visible = False

		chart.category_axis.has_major_tick_marks = False
		chart.category_axis.tick_labels.font.size = fontsize					 
		chart.category_axis.format.line.visible = False
	
		return chart

	def vertical_bar_chart(slide, df, x, y, cx, cy, legend=True, legend_position=XL_LEGEND_POSITION.RIGHT,
							 data_show=True, chart_title=False, title="", fontsize=Pt(10),
							 fontsize_title=Pt(14), percentage=False, bar_width=Pt(10)):
		df.fillna(0, inplace=True)  # Fill NaN values
		# Define chart data
		chart_data = CategoryChartData()
		for i in df.index:
			chart_data.add_category(i)
		for col in df.columns:
			chart_data.add_series(col, np.where(df[col].values == 0, None, df[col].values))
	
	    # Create vertical bar chart
		chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
	
	    # Add legend
		if legend:
			chart.has_legend = True
			chart.legend.include_in_layout = False
			chart.legend.position = legend_position
			chart.legend.font.size = fontsize
		else:
			chart.has_legend = False
	
	    # Show data labels
		if data_show:
			for series_idx, series in enumerate(chart.series):
				for point_idx, point in enumerate(series.points):
					val = df.iloc[point_idx, series_idx]
					label_text = format_number_kmb(val) if val != 0 else ""
					point.data_label.text = label_text
					point.data_label.font.size = fontsize
					point.data_label.position = XL_LABEL_POSITION.OUTSIDE_END
	
	    # Customize category axis (y-axis) label font size
		chart.category_axis.tick_labels.font.size = fontsize
	
	    # Customize value axis (x-axis) label font size and format
		chart.value_axis.tick_labels.font.size = fontsize
		if percentage:
			chart.value_axis.tick_labels.number_format = '0%'
		elif df.max().max() >= 1000:
			chart.value_axis.tick_labels.number_format = '#,##0'  # Add commas for thousands
	
	    # Adjust bar width
		for series in chart.series:
			series.format.line.width = bar_width
	
		# Set chart title
		if chart_title:
			chart.chart_title.text_frame.text = title
			title_font = chart.chart_title.text_frame.paragraphs[0].font
			title_font.bold = True
			title_font.size = fontsize_title
		else:
			chart.has_title = False
	
	    # Remove Gridlines
		value_axis = chart.value_axis
		category_axis = chart.category_axis
		value_axis.has_major_gridlines = False
		value_axis.has_minor_gridlines = False
		category_axis.has_major_gridlines = False
		category_axis.has_minor_gridlines = False
	
		return chart

	def combo_chart(slide, df, x, y, cx, cy, legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                data_show=False, chart_title=False, title="", fontsize=Pt(10), fontsize_title=Pt(12),
                line_width=Pt(1), smooth=False, font_name='Neue Haas Grotesk Text Pro'):

		df.fillna(0, inplace=True)  # Fill NaN values
	
	    # Split Data: Stacked Bar Chart (Categories) & Line Chart (Total)
		df_categories = df.iloc[:, :-1]  # Exclude 'Total' column for bars
		df_total = df.iloc[:, -1]        # Only 'Total' column for line chart
	
	    # Stacked Bar Chart Data
		chart_data = CategoryChartData()
		chart_data.categories = df.index  # Use index as categories (e.g., Months)	
		for category in df_categories.columns:
			chart_data.add_series(category, df_categories[category].values.tolist())
	
	    # Add Stacked Bar Chart
		chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data)
		chart = chart_shape.chart
	
		chart.value_axis.tick_labels.font.size = fontsize
		chart.category_axis.tick_labels.font.size = fontsize
		chart.value_axis.tick_labels.number_format = "#,##0"  # Thousands separator
		chart.value_axis.has_major_gridlines = False
	
	    # Add percentage labels inside bars
		df_percent = df_categories.div(df_total, axis=0) * 100
		for series, category in zip(chart.series, df_percent.columns):
			series.data_labels.show_category_name = False
			series.data_labels.show_value = True
			series.data_labels.font.size = fontsize
			for point, percent in zip(series.points, df_percent[category]):
				point.has_data_label = False
				point.data_label.font.size = fontsize
				point.data_label.text_frame.text = f"{percent:.0f}%"  # Format as percentage
	
	    # Add legend
		if legend:
			chart.has_legend = True
			chart.legend.include_in_layout = False
			chart.legend.position = legend_position
			chart.legend.font.size = fontsize
	
	    # Line Chart Data (Total)
		line_chart_data = CategoryChartData()
		line_chart_data.categories = df.index
		line_chart_data.add_series("Total", df_total.values.tolist())
	
	    # Add Line Chart Overlay
		line_chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, line_chart_data)
		line_chart = line_chart_shape.chart
	
	    # Remove X-axis and Gridlines from Line Chart
		line_chart.category_axis.visible = False # Remove X-axis
		line_chart.value_axis.visible = False # Remove Y-axis
		line_chart.value_axis.has_major_gridlines = False # Remove gridlines
		line_chart.category_axis.has_major_gridlines = False # Remove gridlines
		line_chart.has_legend = False  # Remove legend
		line_chart.has_title = False  # Remove chart title
		line_chart.value_axis.tick_labels.number_format = '""'  # Hide Y-axis labels
	
	
	    # Add Data Labels to Line Chart
		for series in line_chart.series:
			series.has_data_labels = True
			for point in series.points:
				point.has_data_label = True
				point.data_label.number_format = "#,##0"  # Thousands separator
	
	    # Apply Smooth Line if required
		if smooth:
			for series in line_chart.series:
				series.smooth = True
	
	    # Set Line Chart Color and Style
		line_series = line_chart.series[0]
		line_series.marker.style = XL_MARKER_STYLE.CIRCLE
		line_series.marker.format.fill.solid()
		line_series.marker.format.fill.fore_color.rgb = RGBColor(78, 167, 46)
		line_series.format.line.width = line_width
		line_series.format.line.fill.solid()
		line_series.format.line.fill.fore_color.rgb = RGBColor(78, 167, 46)  # Green lime
	
	    # Add Chart Title
		if chart_title:
			chart.chart_title.text_frame.text = title
			chart.chart_title.text_frame.paragraphs[0].font.size = fontsize_title
			chart.chart_title.text_frame.paragraphs[0].font.bold = True
		else:
			chart.has_title = False
		
		return chart

	#----------------
	def combo2_chart(slide, df, x, y, cx, cy, legend=True, legend_position=XL_LEGEND_POSITION.TOP,
			 data_show=False, chart_title=False, title="", fontsize=Pt(10), fontsize_title=Pt(12),
			 line_width=Pt(1), smooth=False, font_name='Neue Haas Grotesk Text Pro'):
		df.fillna(0, inplace=True)

    	# Split data: Bar chart (categories) & Line chart (Total)
		df_categories = df.iloc[:, :-1]  # All but last column
		df_2nd = df.iloc[:, -1]        # Last column (for line)

    	# Bar Chart Data
		chart_data = CategoryChartData()
		chart_data.categories = df.index
		for category in df_categories.columns:
			chart_data.add_series(category, df_categories[category].tolist())

    	# Add Clustered Column Chart (Bar)
		chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
		chart = chart_shape.chart
		chart.value_axis.tick_labels.font.size = fontsize
		chart.category_axis.tick_labels.font.size = fontsize
		chart.value_axis.tick_labels.number_format = "#,##0"
		chart.value_axis.has_major_gridlines = False

    	# Show data labels (if enabled)
		if data_show:
			for series in chart.series:
				series.data_labels.show_value = True
				series.data_labels.font.size = fontsize

    	# Add legend
		if legend:
			chart.has_legend = True
			chart.legend.include_in_layout = False
			chart.legend.position = legend_position
			chart.legend.font.size = fontsize
		else:
			chart.has_legend = False

    	# Line Chart Data
		line_chart_data = CategoryChartData()
		line_chart_data.categories = df.index
		line_chart_data.add_series("value", df_2nd.tolist())

    	# Add Line Chart Overlay
		line_chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, line_chart_data)
		line_chart = line_chart_shape.chart

    	# Hide line chart axes
		line_chart.category_axis.visible = False
		line_chart.value_axis.visible = False
		line_chart.value_axis.has_major_gridlines = False
		line_chart.category_axis.has_major_gridlines = False
		line_chart.has_legend = False
		line_chart.has_title = False
		line_chart.value_axis.tick_labels.number_format = '""'

    	# Add data labels to line chart
		for series in line_chart.series:
			series.has_data_labels = True
			for point in series.points:
				point.has_data_label = True
				point.data_label.number_format = "#,##0"
				point.data_label.font.size = fontsize

    	# Smooth line option
		if smooth:
			for series in line_chart.series:
				series.smooth = True

	# Style line chart
		line_series = line_chart.series[0]
		line_series.marker.style = XL_MARKER_STYLE.CIRCLE
		line_series.marker.format.fill.solid()
		line_series.marker.format.fill.fore_color.rgb = RGBColor(78, 167, 46)  # Lime green
		line_series.format.line.width = line_width
		line_series.format.line.fill.solid()
		line_series.format.line.fill.fore_color.rgb = RGBColor(78, 167, 46)

	# Chart title
		if chart_title:
			chart.chart_title.text_frame.text = title
			p = chart.chart_title.text_frame.paragraphs[0]
			p.font.size = fontsize_title
			p.font.bold = True
		else:
			chart.has_title = False

		return chart
	#----------------

	def adjust_dataframe(df, columns, index=False):
		# Combine existing and desired columns/index
		all_columns = list(df.columns) + [col for col in columns if col not in df.columns]
		if index == False:
			df = df.reindex(columns=all_columns,fill_value=np.nan)
		else:
			all_index = list(df.index) + [idx for idx in index if idx not in df.index]
			# Reindex to include all columns and index values
			df = df.reindex(columns=all_columns, index=all_index, fill_value=np.nan)
	
		return df
	#-----------------
	def df_to_bullets(slide, df, left=Inches(1), top=Inches(1), width=Inches(5), height=Inches(4), font_size=9, font_bold = False):

	# Add textbox
		textbox = slide.shapes.add_textbox(left, top, width, height)
		text_frame = textbox.text_frame
		text_frame.word_wrap = True
		text_frame.clear()  # Clear default paragraph

		for value in df.iloc[:, 0]:  # Take values from the first (only) column
			p = text_frame.add_paragraph()
			p.text = str(value)
			p.level = 0
			p.font.size = Pt(font_size)
			p.font.color.rgb = RGBColor(0, 0, 0)
			p.alignment = PP_ALIGN.CENTER
	
	#----------------
	
	def df_to_labeled_bullets(slide, df, left=Inches(1), top=Inches(1), width=Inches(5), height=Inches(4), font_size=9, font_bold = False):
	# Add textbox
		textbox = slide.shapes.add_textbox(left, top, width, height)
		text_frame = textbox.text_frame
		text_frame.word_wrap = True
		text_frame.clear()  # Clear default paragraph
	
	# Combine label (first column) and value (second column)
		for i in range(len(df)):
			label = str(df.iloc[i, 0])
			value = str(df.iloc[i, 1])
			p = text_frame.add_paragraph()
			p.text = f"{label}: {value}"
			p.level = 0
			p.font.size = Pt(font_size)
			p.font.color.rgb = RGBColor(0, 0, 0)
			p.font.bold = font_bold
			p.alignment = PP_ALIGN.CENTER

	# ---------------
	# --- Get thumbnails
	api_key = st.secrets["rapidapi_key"]["rapidapi_key"]

	def get_tiktok_thumbnail(video_url, api_key):
		url = "https://tiktok-video-downloader-api.p.rapidapi.com/media-info/"
		querystring = {"url": video_url}
		headers = {
			"X-RapidAPI-Key": api_key,
			"X-RapidAPI-Host": "tiktok-video-downloader-api.p.rapidapi.com"
		}
		response = requests.get(url, headers=headers, params=querystring)
		if response.status_code == 200:
			data = response.json()
			thumbnail_url = data.get("thumbnail_url")
			return thumbnail_url
		else:
			print(f"Error: {response.status_code}")
			return None

	def get_instagram_thumbnail(instagram_url):
		try:
			headers = {'User-Agent': 'Mozilla/5.0'}
			response = requests.get(instagram_url, headers=headers)
			if response.status_code == 200:
				soup = BeautifulSoup(response.text, 'html.parser')
				og_image = soup.find('meta', property='og:image')
				if og_image:
					return og_image['content']
				return None
		except:
			return None

	def get_thumbnail(link, api_key):
		if not link:
			return None
		if "tiktok.com" in link:
			return get_tiktok_thumbnail(link, api_key)
		elif "instagram.com" in link:
			return get_instagram_thumbnail(link)
		else:
			return None
			
	def download_image_from_url(image_url, save_path):
		try:
			response = requests.get(image_url)
			response.raise_for_status()
			with open(save_path, 'wb') as f:
				f.write(response.content)
			return save_path
		except Exception as e:
			print(f"Error downloading image from {image_url}: {e}")
			return None
			
	#-----------------

	# Load credentials from Streamlit Secrets
	credentials_dict = st.secrets["gcp_service_account"]
	credentials = service_account.Credentials.from_service_account_info(credentials_dict)

	# Initialize BigQuery Client
	client = bigquery.Client(credentials=credentials, project=credentials.project_id)
   
# df for Competitive data (TDK)
	# Query Code
	query = """
 	SELECT 
	date
	,month
	,years
	,brand
	,category
	,division
	,manufacturer
	,advertiser_name
	,SUM(views_float) as views
	,SUM(engagements) as engagements
	,SUM(content) as content
	FROM loreal-id-prod.loreal_storage.advocacy_tdk_df
	WHERE years = {} or years = {}
	GROUP BY 
	date
	,month
	,years
	,brand
	,division
	,category
	,manufacturer
	,advertiser_name
	""".format(year_map, (year_map-1))

	# Fetch data
	df = client.query(query).to_dataframe()

	# Add Date & Quarter column
	df['date'] = pd.to_datetime(df['date'], format='%y-%m-%d')

	df['quarter'] = (df['date'].dt.quarter).map({1: 'Q1', 2: 'Q2', 3: 'Q3', 4: 'Q4'})

# df2 for INTERNAL DATA
	# Query Code 2
	query2 = f"""
	SELECT 
 		DATE(date_post) AS date,
  		EXTRACT(MONTH FROM date_post) AS month,
   		EXTRACT(YEAR FROM date_post) AS years,
    		brand,
		sub_brand,
		division,
		sub_category AS category,
		tier,
  		kol_name,
    		kol_persona,
      		followers,
  		link_post,
		campaign,
		spark_ads AS advocacy,
		rate,
		actual_views as views,
		engagement as engagements,
		er_content,
		1 AS content
  		FROM loreal-id-prod.loreal_storage.advocacy_campaign_df
    		WHERE EXTRACT(YEAR FROM date_post) = {year}
  	"""

	# Fetch data
	df2 = client.query(query2).to_dataframe()

	# Add Date & Quarter column
	df2['date'] = pd.to_datetime(df2['date'], format='%y-%m-%d')

	df2['quarter'] = (df2['date'].dt.quarter).map({1: 'Q1', 2: 'Q2', 3: 'Q3', 4: 'Q4'})

	df2['years'] = df2['years'].astype(str)
	df2['month'] = df2['month'].astype(str)
	
#---- SLIDES PRESENTATION ----
#	ppt_temp_loc = "GOAT_LOreal_DeckToGo/pages/Template Deck to Go - Loreal Indonesia.pptx"
  
	ppt = Presentation(uploaded_file)
		
#------------PAGE 1--------------
	page_no = 0 #PAGE1

## Title Slide
# Add a title to the slide
	format_title(ppt.slides[page_no], "MONTHLY BRAND & DIVISION REPORT", alignment=PP_ALIGN.LEFT, font_name= 'Aptos Display', font_size=50, font_bold=True,left=Pt(35), top=Pt(150), width=Pt(500), height=Pt(70), font_color=RGBColor(255, 255, 255))

	st.write("Slide 1 of 17")

#------------PAGE 2--------------
	page_no = page_no + 1 #PAGE2

	# Add a title to the slide
	format_title(ppt.slides[page_no], "MONTHLY SOV & SOE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], "selected month: " + str(month), alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=10, font_bold=False,left=Inches(0.5), top=Inches(1.5), width=Inches(12), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

	format_title(ppt.slides[page_no], "Total Views", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(5), top=Inches(2), width=Inches(1.43), height=Inches(1.01), font_color=RGBColor(255, 255, 255))
	format_title(ppt.slides[page_no], "Total Eng.", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(11), top=Inches(2), width=Inches(1.43), height=Inches(1.01), font_color=RGBColor(255, 255, 255))

	# Filter the dataframe
	df_m = df[(df['division'].isin(division)) & (df['category'].isin(category)) & (df['years'] == year_map) &  (df['month'].isin(month))]

	# Perform groupby and aggregation with handling for datetime64 columns
	grouped_df_m = df_m.groupby('brand').agg({
    	# Numerical columns: sum them
    	'views': 'sum',
    	'engagements': 'sum',
    	'content': 'sum'
	})
	grouped_df_m = grouped_df_m.reset_index()

	# Calculate Total Views & Total Engagements
	total_views_m = grouped_df_m['views'].sum().astype(int)
	total_engagement_m = grouped_df_m['engagements'].sum().astype(int)

	# Calculate SOV (%)
	grouped_df_m['SOV%'] = (grouped_df_m['views'] / total_views_m)
	sov_df_m = grouped_df_m[['brand', 'SOV%']].sort_values(by='SOV%', ascending=False)

	# Sort by 'SOV%' and keep the top 14 brands
	top_brands_sov = sov_df_m.head(14)

	# Group all other brands into "Others"
	others_sov = 1 - top_brands_sov["SOV%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_sov = top_brands_sov.copy()
	if others_sov > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOV%": others_sov}])
		top_brands_sov = pd.concat([top_brands_sov, others_row], ignore_index=True)

	# Calculate SOE (%)
	grouped_df_m['SOE%'] = (grouped_df_m['engagements'] / total_engagement_m)
	soe_df_m = grouped_df_m[['brand', 'SOE%']].sort_values(by='SOE%', ascending=False)

	# Sort by 'SOE%' and keep the top 14 brands
	top_brands_soe = soe_df_m.head(14)

	# Group all other brands into "Others"
	others_soe = 1 - top_brands_soe["SOE%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_soe = top_brands_soe.copy()
	if others_soe > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOE%": others_soe}])
		top_brands_soe = pd.concat([top_brands_soe, others_row], ignore_index=True)

	
	# Add pie chart
	
	pie_chart(ppt.slides[page_no], top_brands_sov.set_index('brand'), Inches(0.5), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOV', fontsize_title = Pt(20), fontsize=9)
	pie_chart(ppt.slides[page_no], top_brands_soe.set_index('brand'), Inches(7), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOE', fontsize_title = Pt(20), fontsize=9)

	format_title(ppt.slides[page_no], "Total Views", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_italic=True,left=Inches(5.3), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_views_m, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(5.3), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))

	format_title(ppt.slides[page_no], "Total Eng.", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_italic=True,left=Inches(11.7), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_engagement_m, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(11.7), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))

	st.write("Slide 2 of 17")
	
#------------PAGE 3--------------
	page_no = page_no + 1 #PAGE3

# Add a title to the slide
	format_title(ppt.slides[page_no], "QUARTERLY SOV & SOE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], "selected quarter: " + str(quarter), alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=10, font_bold=False,left=Inches(0.5), top=Inches(1.5), width=Inches(12), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Filter the dataframe
	df_q = df[(df['division'].isin(division)) & (df['category'].isin(category)) & (df['years'] == year_map) & (df['quarter'].isin(quarter))]

# Perform groupby and aggregation with handling for datetime64 columns
	grouped_df_q = df_q.groupby('brand').agg({
    	# Numerical columns: sum them
    	'views': 'sum',
    	'engagements': 'sum',
    	'content': 'sum'
	})
	grouped_df_q = grouped_df_q.reset_index()

# Calculate Total Views
	total_views_q = (grouped_df_q['views'].sum()).astype(int)
	total_engagement_q = (grouped_df_q['engagements'].sum()).astype(int)

	# Calculate SOV (%)
	grouped_df_q['SOV%'] = (grouped_df_q['views'] / total_views_q)
	sov_df_q = grouped_df_q[['brand', 'SOV%']].sort_values(by='SOV%', ascending=False)

	# Sort by 'SOV%' and keep the top 14 brands
	top_brands_sov = sov_df_q.head(14)

	# Group all other brands into "Others"
	others_sov = 1 - top_brands_sov["SOV%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_sov = top_brands_sov.copy()
	if others_sov > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOV%": others_sov}])
		top_brands_sov = pd.concat([top_brands_sov, others_row], ignore_index=True)

	# Calculate SOE (%)
	grouped_df_q['SOE%'] = (grouped_df_q['engagements'] / total_engagement_q)
	soe_df_q = grouped_df_q[['brand', 'SOE%']].sort_values(by='SOE%', ascending=False)

	# Sort by 'SOE%' and keep the top 14 brands
	top_brands_soe = soe_df_q.head(14)

	# Group all other brands into "Others"
	others_soe = 1 - top_brands_soe["SOE%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_soe = top_brands_soe.copy()
	if others_soe > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOE%": others_soe}])
		top_brands_soe = pd.concat([top_brands_soe, others_row], ignore_index=True)

	
	# Add pie chart
	
	pie_chart(ppt.slides[page_no], top_brands_sov.set_index('brand'), Inches(0.5), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOV', fontsize_title = Pt(20), fontsize=9)
	pie_chart(ppt.slides[page_no], top_brands_soe.set_index('brand'), Inches(7), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOE', fontsize_title = Pt(20), fontsize=9)

	format_title(ppt.slides[page_no], "Total Views", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(5.3), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_views_q, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(5.3), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))

	format_title(ppt.slides[page_no], "Total Eng.", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(11.7), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_engagement_q, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(11.7), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))
	
	st.write("Slide 3 of 17")
	
#------------PAGE 4--------------
	page_no = page_no + 1 #PAGE4

# Add a title to the slide
	format_title(ppt.slides[page_no], "ANNUAL SOV & SOE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Filter the dataframe
	df_y = df[(df['division'].isin(division)) & (df['category'].isin(category)) & (df['years'] == year_map)]

# Perform groupby and aggregation with handling for datetime64 columns
	grouped_df_y = df_y.groupby('brand').agg({
		# Numerical columns: sum them
		'views': 'sum',
		'engagements': 'sum',
		'content': 'sum'
	})
	grouped_df_y = grouped_df_y.reset_index()

# Calculate Total Views, Engagements & Content
	total_views_y = (grouped_df_y['views'].sum()).astype(int)
	total_engagement_y = (grouped_df_y['engagements'].sum()).astype(int)
	total_content_y = (grouped_df_y['content'].sum()).astype(int)

# Calculate SOV (%)
	grouped_df_y['SOV%'] = (grouped_df_y['views'] / total_views_y)
	sov_df_y = grouped_df_y[['brand', 'SOV%']].sort_values(by='SOV%', ascending=False)

	# Sort by 'SOV%' and keep the top 14 brands
	top_brands_sov = sov_df_y.head(14)

	# Group all other brands into "Others"
	others_sov = 1 - top_brands_sov["SOV%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_sov = top_brands_sov.copy()
	if others_sov > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOV%": others_sov}])
		top_brands_sov = pd.concat([top_brands_sov, others_row], ignore_index=True)

	# Calculate SOE (%)
	grouped_df_y['SOE%'] = (grouped_df_y['engagements'] / total_engagement_y)
	soe_df_y = grouped_df_y[['brand', 'SOE%']].sort_values(by='SOE%', ascending=False)

	# Sort by 'SOE%' and keep the top 14 brands
	top_brands_soe = soe_df_y.head(14)

	# Group all other brands into "Others"
	others_soe = 1 - top_brands_soe["SOE%"].sum()  # Remaining percentage

	# Append "Others" row if applicable
	top_brands_soe = top_brands_soe.copy()
	if others_soe > 0:
		others_row = pd.DataFrame([{"brand": "Others", "SOE%": others_soe}])
		top_brands_soe = pd.concat([top_brands_soe, others_row], ignore_index=True)
# Add pie chart
	
	pie_chart(ppt.slides[page_no], top_brands_sov.set_index('brand'), Inches(0.5), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOV', fontsize_title = Pt(20), fontsize=9)
	pie_chart(ppt.slides[page_no], top_brands_soe.set_index('brand'), Inches(7), Inches(1.5), Inches(6), Inches(6), chart_title=True, title='SOE', fontsize_title = Pt(20), fontsize=9)

	format_title(ppt.slides[page_no], "Total Views", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(5.3), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_views_y, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(5.3), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))

	format_title(ppt.slides[page_no], "Total Eng.", alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_italic=True,left=Inches(11.7), top=Inches(3), width=Inches(1.3), height=Inches(1.01), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format(total_engagement_y, ","), alignment=PP_ALIGN.CENTER, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(11.7), top=Inches(2.7), width=Inches(1.3), height=Inches(0.5), font_color=RGBColor(0, 0, 0))

	st.write("Slide 4 of 17")
	
#------------PAGE 5--------------
	page_no = page_no + 1 #PAGE5

# Add a title to the slide
	format_title(ppt.slides[page_no], "YTD BEAUTY MARKET VIEWS AND ENGAGEMENT SUMMARY", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.3), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Aggregate Views by Month and Manufacturer
	stacked_data_views = df_y.pivot_table(index="month", columns="manufacturer", values="views", aggfunc="sum", fill_value=0)
	stacked_data_views['Total'] = stacked_data_views.sum(axis=1)

# Aggregate Views by Month and Manufacturer
	stacked_data_eng = df_y.pivot_table(index="month", columns="manufacturer", values="engagements", aggfunc="sum", fill_value=0)
	stacked_data_eng['Total'] = stacked_data_eng.sum(axis=1)

# Sort months correctly
	month_order = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
	stacked_data_eng = stacked_data_eng.reindex(month_order)
	stacked_data_views = stacked_data_views.reindex(month_order)

# Add combo stacked bar chart
	combo_chart(ppt.slides[page_no], stacked_data_views, Inches(.1), Inches(1.5), Inches(9), Inches(2.8), chart_title=True, title="Market Movement - Views",
            fontsize=Pt(10), fontsize_title=Pt(12), smooth=True, data_show=True)
	combo_chart(ppt.slides[page_no], stacked_data_eng, Inches(.1), Inches(4.5), Inches(9), Inches(2.8), chart_title=True, title="Market Movement - Engagements",
            fontsize=Pt(10), fontsize_title=Pt(12), smooth=True)
	
	st.write("Slide 5 of 17")
	
#------------PAGE 6--------------
	page_no = page_no + 1 #PAGE6

	format_title(ppt.slides[page_no], "VIEWS PERFORMANCE: YTD and YTG with SOV PROJECTION", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Filter the dataframe
	df_y2 = df[(df['division'].isin(division)) & (df['years'] == year_map)]

# Calculate Total Views per Category
	total_views = df_y2.groupby('category')['views'].sum().reset_index()
	total_views.rename(columns={'views': 'Total_Views'}, inplace=True)

# Aggregate Views and calculate SOV
	df_grouped = df_y2.groupby(['brand', 'category', 'advertiser_name'])['views'].sum().reset_index()
	df_grouped = df_grouped.merge(total_views, on='category', how='left')
	df_grouped['SOV'] = (df_grouped['views'] / df_grouped['Total_Views']).map('{:.0%}'.format)
	df_grouped['views'] = df_grouped['views'].astype(int)

# Rank Brands within each Category
	df_grouped['#Rank'] = df_grouped.groupby('category')['views'].rank(method='dense', ascending=False).astype(int)

# Filter for L'Oreal Advertiser
	df_final = df_grouped[df_grouped['advertiser_name'] == "L'Oreal"][['category', 'brand', 'views', 'SOV', '#Rank']]
	df_final = df_final.sort_values(['category', 'brand'])

# Add table
	table_default(ppt.slides[page_no], df_final, Inches(1), Inches(1.2), Inches(12.2), Inches(5.2),
				  [Inches(1.5)]*2+[Inches(0.75)]*3+[Inches(1),Inches(0.5)], Inches(0.5), header=True, upper=True, fontsize=12, alignment=PP_ALIGN.LEFT)
	
	st.write("Slide 6 of 17")
	
#------------PAGE 7--------------
	page_no = page_no + 1 #PAGE7

	format_title(ppt.slides[page_no], "MONTHLY TIMELINE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Aggregate total views per brand and rank them
	brand_ranking = df_y.groupby('brand')['views'].sum().reset_index()
	brand_ranking['Brand_Rank'] = brand_ranking['views'].rank(method='dense', ascending=False)

# Keep only the top 7 brands
	top_brands = brand_ranking[brand_ranking['Brand_Rank'] <= 7]['brand']

# Retrieve daily views for these top 7 brands
	df_top_brands = df_y[df_y['brand'].isin(top_brands)]
	df_m_views = pd.pivot_table(df_top_brands[['date','brand','views']], columns = 'date', index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_m_views.columns = df_m_views.columns.droplevel() # drop column level
	df_m_views.columns = df_m_views.columns.strftime('%b')

# Reshape the DataFrame before grouping to have 'views' as a column
	df_m_views2 = df_m_views.reset_index().melt(id_vars=['brand'], var_name='month', value_name='views')

	df_views = df_m_views2.groupby(['brand'], as_index=False)['views'].sum()
	df_views['SOV%'] = (df_views['views'] / total_views_y).map('{:.1%}'.format)
	df_views['views'] = df_views['views'].astype(int)
	df_views['rank'] = df_views['views'].rank(method='dense', ascending=False).astype(int)
	df_views = df_views.sort_values('rank')

# Add line chart
	line_marker_chart(ppt.slides[page_no], df_m_views, Inches(0.5), Inches(1.7), Inches(7), Inches(5), legend=True, legend_position=XL_LEGEND_POSITION.BOTTOM, chart_title = False, fontsize_title = Pt(12), line_width = Pt(2), smooth=True)
# Add table
	table_default(ppt.slides[page_no], df_views, Inches(8), Inches(1.7), Inches(5), Inches(5),
								[Inches(1.5)]+[Inches(0.75)]*3+[Inches(1),Inches(0.5)], Inches(0.5), header=True, upper=True, fontsize=12, alignment=PP_ALIGN.LEFT)
	st.write("Slide 7 of 17")
	
#------------PAGE 8--------------
	page_no = page_no + 1 
	format_title(ppt.slides[page_no], "QUARTERLY TIMELINE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

	# Aggregate total views per brand and rank them
	df_q_views = pd.pivot_table(df_top_brands[['quarter','brand','views']], columns = 'quarter', index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_q_views.columns = df_q_views.columns.droplevel() # drop column level

# Add line chart
	line_marker_chart(ppt.slides[page_no], df_q_views, Inches(0.5), Inches(1.7), Inches(7), Inches(5), legend=True, legend_position=XL_LEGEND_POSITION.BOTTOM, chart_title = False, fontsize_title = Pt(12), line_width = Pt(2), smooth=True)
# Add table
	table_default(ppt.slides[page_no], df_views, Inches(8), Inches(1.7), Inches(5), Inches(5),
								[Inches(1.5)]+[Inches(0.75)]*3+[Inches(1),Inches(0.5)], Inches(0.5), header=True, upper=True, fontsize=12, alignment=PP_ALIGN.LEFT)

	st.write("Slide 8 of 17")
	
#------------PAGE 9--------------
	page_no = page_no + 1 
	format_title(ppt.slides[page_no], "BAR CHART COMPARISON QTR VERSUS", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

	if isinstance(quarter, list) and len(quarter) == 2:
		quarter_compare = quarter
	else:
		_quarter = quarter[0]
		quarter_mapping = {'Q1': ['Q4', 'Q1'], 'Q2': ['Q1', 'Q2'], 'Q3': ['Q2', 'Q3'], 'Q4': ['Q3', 'Q4']}
		quarter_compare = quarter_mapping.get(_quarter, [_quarter])
	title_q = quarter_compare[0] + " vs " + quarter_compare[1]

	if quarter_ == ["Quarter 1"]:
		df_y_qq = df[(df['division'].isin(division)) & (df['category'].isin(category)) & (df['years'].isin(year_range))]
	else:
		df_y_qq = df[(df['division'].isin(division)) & (df['category'].isin(category)) & (df['years'] == year_map)]
	
	df_top_brands = df_y_qq[df_y_qq['brand'].isin(top_brands)]  

	df_q_views = pd.pivot_table(df_top_brands[['quarter','brand','views']], columns = 'quarter', index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_q_content = pd.pivot_table(df_top_brands[['quarter','brand','content']], columns = 'quarter', index = 'brand', aggfunc = 'sum', fill_value = 0)

	df_q_views = df_q_views[[('views', q) for q in quarter_compare]]
	df_q_views = np.ceil(df_q_views * 10) / 10 # Round up with 1 decimal place
	df_q_views.columns = quarter_compare

	df_q_content = df_q_content[[('content', q) for q in quarter_compare]]
	df_q_content.columns = quarter_compare

# Add horizontal bar chart views
	horizontal_bar_chart(ppt.slides[page_no], df_q_views, Inches(0.5), Inches(1.9), Inches(5.5), Inches(5),
                     chart_title = True, title= f"{title_q} Views Contribution", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))

# Add horizontal bar chart content
	horizontal_bar_chart(ppt.slides[page_no], df_q_content, Inches(7), Inches(1.9), Inches(5.5), Inches(5),
                     chart_title = True, title= f"{title_q} Content Contribution", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	
	st.write("Slide 9 of 17")
	
#------------PAGE 10--------------
	page_no = page_no + 1
	format_title(ppt.slides[page_no], "MONTHLY PER BRAND CONTRIBUTION", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Filter the dataframe
	df_10 = df2[(df2['brand'].isin(brands)) & (df2['years'] == year) &  (df2['month'].isin(month_num))]
	df_10_views = pd.pivot_table(df_10[['brand', 'views']], index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_10_eng = pd.pivot_table(df_10[['brand', 'engagements']], index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_10_content = pd.pivot_table(df_10[['brand', 'content']], index = 'brand', aggfunc = 'sum', fill_value = 0)

# Add vertical bar chart
	vertical_bar_chart(ppt.slides[page_no], df_10_views, Inches(0.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "View Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	vertical_bar_chart(ppt.slides[page_no], df_10_eng, Inches(4.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "Engagement Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	vertical_bar_chart(ppt.slides[page_no], df_10_content, Inches(8.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "Content Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	
# Add rectangle
	shapes1 = ppt.slides[page_no].shapes
	rectangle1 = shapes1.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5), Inches(3.7), Inches(2))
	rectangle1.adjustments[0] = 0.1
	rectangle1.fill.background()
	text_frame1 = rectangle1.text_frame
	text_frame1.text = "1. ..." 
	text_frame1.paragraphs[0].font.size = Pt(13)

	shapes2 = ppt.slides[page_no].shapes
	rectangle2 = shapes2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(5), Inches(3.7), Inches(2))
	rectangle2.adjustments[0] = 0.1
	rectangle2.fill.background()
	text_frame2 = rectangle2.text_frame
	text_frame2.text = "1. ..."
	text_frame2.paragraphs[0].font.size = Pt(13)

	shapes3 = ppt.slides[page_no].shapes
	rectangle3 = shapes3.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(5), Inches(3.7), Inches(2))
	rectangle3.adjustments[0] = 0.1
	rectangle3.fill.background()
	text_frame3 = rectangle3.text_frame
	text_frame3.text = "1. ..."
	text_frame3.paragraphs[0].font.size = Pt(13)

	st.write("Slide 10 of 17")
	
#------------PAGE 11--------------
	page_no = page_no + 1 
	format_title(ppt.slides[page_no], "QUARTERLY PER BRAND CONTRIBUTION", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	
	df_11 = df2[(df2['brand'].isin(brands)) & (df2['years'] == year) &  (df2['quarter'].isin(quarter))]
	df_11_views = pd.pivot_table(df_11[['brand', 'views']], index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_11_eng = pd.pivot_table(df_11[['brand', 'engagements']], index = 'brand', aggfunc = 'sum', fill_value = 0)
	df_11_content = pd.pivot_table(df_11[['brand', 'content']], index = 'brand', aggfunc = 'sum', fill_value = 0)

# Add vertical bar chart
	vertical_bar_chart(ppt.slides[page_no], df_11_views, Inches(0.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "View Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	vertical_bar_chart(ppt.slides[page_no], df_11_eng, Inches(4.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "Engagement Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	vertical_bar_chart(ppt.slides[page_no], df_11_content, Inches(8.9), Inches(1.9), Inches(3.7), Inches(3),
                     chart_title = True, title= "Content Performance", fontsize_title = Pt(16),
                     legend=True, legend_position=XL_LEGEND_POSITION.TOP,
                     bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	
# Add rectangle
	shapes1 = ppt.slides[page_no].shapes
	rectangle1 = shapes1.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5), Inches(3.7), Inches(2))
	rectangle1.adjustments[0] = 0.1
	rectangle1.fill.background()
	text_frame1 = rectangle1.text_frame
	text_frame1.text = "1. ..." 
	text_frame1.paragraphs[0].font.size = Pt(13)

	shapes2 = ppt.slides[page_no].shapes
	rectangle2 = shapes2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(5), Inches(3.7), Inches(2))
	rectangle2.adjustments[0] = 0.1
	rectangle2.fill.background()
	text_frame2 = rectangle2.text_frame
	text_frame2.text = "1. ..."
	text_frame2.paragraphs[0].font.size = Pt(13)

	shapes3 = ppt.slides[page_no].shapes
	rectangle3 = shapes3.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(5), Inches(3.7), Inches(2))
	rectangle3.adjustments[0] = 0.1
	rectangle3.fill.background()
	text_frame3 = rectangle3.text_frame
	text_frame3.text = "1. ..."
	text_frame3.paragraphs[0].font.size = Pt(13)
	
	st.write("Slide 11 of 17")
	
#------------PAGE 12--------------
	page_no = page_no + 1
	format_title(ppt.slides[page_no], "BRAND SCORE-CARD by CATEGORY", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	
	df_12 = df2[(df2['brand'].isin(brands)) & (df2['years'] == year)]
	df_12 = pd.pivot_table(df_12[['category', 'brand', 'sub_brand', 'rate', 'views', 'engagements', 'content']], index= ['category', 'brand', 'sub_brand'],
			       values = ['rate', 'views', 'engagements', 'content'], aggfunc = 'sum', fill_value = 0)
	df_12 = df_12.reset_index()
	df_12['rate'] = df_12['rate'].fillna(0).astype(int)
	df_12['eng_rate'] = np.where(df_12['views'] != 0, df_12['engagements'] / df_12['views'], 0)
	df_12['eng_rate'] = df_12['eng_rate'].map('{:.1%}'.format)
	df_12['CPV'] =  np.where(df_12['views'] != 0, df_12['rate'] / df_12['views'], 0)
	df_12['CPV'] = np.ceil(df_12['CPV'] * 100) / 100
	df_12['CPE'] = np.where(df_12['engagements'] != 0, df_12['rate'] / df_12['engagements'], 0)
	df_12['CPE'] = np.ceil(df_12['CPE'] * 100) / 100
	
	df_12['category'] = df_12['category'].where(
	~df_12['category'].duplicated(keep='first'), df_12['category'] + '_' + df_12.groupby('category').cumcount().astype(str))
	df_12 = df_12[['category', 'brand', 'sub_brand', 'rate', 'views', 'engagements', 'content',  'CPV', 'eng_rate', 'CPE']]
	df_12_transpose = df_12.transpose()
	df_12_transpose.reset_index(inplace=True)
	
	df_12_transpose.columns = df_12_transpose.iloc[0]
	df_12_transpose = df_12_transpose[1:].reset_index(drop=True)
	df_12_transpose.rename(columns={'category': 'Market'}, inplace=True)
	num_columns = df_12_transpose.shape[1]
	
# Add table	
	table_default(ppt.slides[page_no], df_12_transpose, Inches(0.5), Inches(1.2), Inches(12.2), Inches(7),
		      [Inches(1)]*num_columns, Inches(0.5), header=True, upper=True, fontsize=10, alignment=PP_ALIGN.LEFT)
	
	st.write("Slide 12 of 17")

#------------PAGE 13--------------
	page_no = page_no + 1
	format_title(ppt.slides[page_no], "BRAND SCORE-CARD by DIVISION", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	
	df_13= df2[(df2['brand'].isin(brands)) & (df2['years'] == year)]
	df_13 = pd.pivot_table(df_13[['division', 'brand', 'sub_brand', 'rate', 'views', 'engagements', 'content']], index= ['division', 'brand', 'sub_brand'],
			       values = ['rate', 'views', 'engagements', 'content'], aggfunc = 'sum', fill_value = 0)
	df_13 = df_13.reset_index()
	df_13['rate'] = df_13['rate'].fillna(0).astype(int)
	df_13['eng_rate'] = np.where(df_13['views'] != 0, df_13['engagements'] / df_13['views'], 0)
	df_13['eng_rate'] = df_13['eng_rate'].map('{:.1%}'.format)
	df_13['CPV'] =  np.where(df_13['views'] != 0, df_13['rate'] / df_13['views'], 0)
	df_13['CPV'] = np.ceil(df_13['CPV'] * 100) / 100
	df_13['CPE'] = np.where(df_13['engagements'] != 0, df_13['rate'] / df_13['engagements'], 0)
	df_13['CPE'] = np.ceil(df_13['CPE'] * 100) / 100
	
	df_13['division'] = df_13['division'].where(
	~df_13['division'].duplicated(keep='first'), df_13['division'] + '_' + df_13.groupby('division').cumcount().astype(str))
	df_13 = df_13[['division', 'brand', 'sub_brand', 'rate', 'views', 'engagements', 'content',  'CPV', 'eng_rate', 'CPE']]
	df_13_transpose = df_13.transpose()
	df_13_transpose.reset_index(inplace=True)
	
	df_13_transpose.columns = df_13_transpose.iloc[0]
	df_13_transpose = df_13_transpose[1:].reset_index(drop=True)
	df_13_transpose.rename(columns={'division': 'Division'}, inplace=True)
	num_columns = df_13_transpose.shape[1]
	
# Add table	
	table_default(ppt.slides[page_no], df_13_transpose, Inches(0.5), Inches(1.2), Inches(12.2), Inches(7),
		      [Inches(1)]*num_columns, Inches(0.5), header=True, upper=True, fontsize=10, alignment=PP_ALIGN.LEFT)
	
	st.write("Slide 13 of 17")	
	
#------------PAGE 14--------------
	page_no = page_no + 1
	format_title(ppt.slides[page_no], "---TITLE---", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	
	df_14 = df_y # data TDK YTD
	
# Aggregate total views per brand and rank them
	rank_view = df_14.groupby('brand')['views'].sum().reset_index()
	rank_view['views'] = np.ceil(rank_view['views'] * 10) / 10
	rank_view['Rank'] = rank_view['views'].rank(method='dense', ascending=False)
	rank_view['SOV%'] = (rank_view['views'] / total_views_y).map('{:.1%}'.format)
	rank_view.sort_values('Rank', ascending=True, inplace=True)

# Keep only the top 10 brands
	top10_views_brands = rank_view[rank_view['Rank'] <= 10]['brand']

# Aggregate total engagements per brand and rank them
	rank_eng = df_14.groupby('brand')['engagements'].sum().reset_index()
	rank_eng['engagements'] = np.ceil(rank_eng['engagements'] * 10) / 10
	rank_eng['Rank'] = rank_eng['engagements'].rank(method='dense', ascending=False)
	rank_eng['SOE%'] = (rank_eng['engagements'] / total_engagement_y).map('{:.1%}'.format)
	rank_eng.sort_values('Rank', ascending=True, inplace=True)

# Keep only the top 10 brands
	top10_eng_brands = rank_eng[rank_eng['Rank'] <= 10]['brand']

# Aggregate total content per brand and rank them
	rank_content = df_14.groupby('brand')['content'].sum().reset_index()
	rank_content['content'] = np.ceil(rank_content['content'] * 10) / 10
	rank_content['Rank'] = rank_content['content'].rank(method='dense', ascending=False)
	rank_content['SOC%'] = (rank_content['content'] / total_content_y).map('{:.1%}'.format)
	rank_content.sort_values('Rank', ascending=True, inplace=True)

# Keep only the top 10 brands
	top10_content_brands = rank_content[rank_content['Rank'] <= 10]['brand']

	df_14_views = pd.pivot_table(df_14[(df_14['brand'].isin(top10_views_brands))], index = 'brand', values= 'views', aggfunc = 'sum', fill_value = 0).sort_values('views', ascending=True)
	df_14_views['views'] = np.ceil(df_14_views['views'] * 10) / 10 # Round up with 1 decimal place
	
	df_14_eng = pd.pivot_table(df_14[(df_14['brand'].isin(top10_eng_brands))], index = 'brand', values= 'engagements', aggfunc = 'sum', fill_value = 0).sort_values('engagements', ascending=True)
	df_14_eng['engagements'] = np.ceil(df_14_eng['engagements'] * 10) / 10 # Round up with 1 decimal place

	df_14_content = pd.pivot_table(df_14[(df_14['brand'].isin(top10_content_brands))], index = 'brand', values= 'content', aggfunc = 'sum', fill_value = 0).sort_values('content', ascending=True)
	df_14_content['content'] = np.ceil(df_14_content['content'] * 10) / 10 # Round up with 1 decimal place

	# Add horizontal bar chart views
	horizontal_bar_chart(ppt.slides[page_no], df_14_views, Inches(0.5), Inches(1.9), Inches(3.5), Inches(5),
                     chart_title = True, title= "SOV", fontsize_title = Pt(16),
                     legend=False, bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	horizontal_bar_chart(ppt.slides[page_no], df_14_eng, Inches(4.5), Inches(1.9), Inches(3.5), Inches(5),
                     chart_title = True, title= "SOE", fontsize_title = Pt(16),
                     legend=False, bar_width = Pt(8), percentage=False, fontsize=Pt(10))
	horizontal_bar_chart(ppt.slides[page_no], df_14_content, Inches(8.5), Inches(1.9), Inches(3.5), Inches(5),
                     chart_title = True, title= "SOC", fontsize_title = Pt(16),
                     legend=False, bar_width = Pt(8), percentage=False, fontsize=Pt(10))

	# Add Total Views, Total Eng, Total Content
	format_title(ppt.slides[page_no], "Views", alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=10.5,left=Inches(4), top=Inches(6.5), width=Inches(0.8), height=Inches(0.6), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format_number_kmb(total_views_y), alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=14, font_bold=True,left=Inches(4), top=Inches(6.2), width=Inches(0.8), height=Inches(0.6), font_color=RGBColor(146, 39, 143))

	format_title(ppt.slides[page_no], "Engagements", alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=10.5, left=Inches(7.8), top=Inches(6.5), width=Inches(1.3), height=Inches(0.6), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format_number_kmb(total_engagement_y), alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=14, font_bold=True,left=Inches(7.8), top=Inches(6.2), width=Inches(1.3), height=Inches(0.6), font_color=RGBColor(146, 39, 143))

	format_title(ppt.slides[page_no], "Contents", alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=10.5, font_italic=False,left=Inches(12), top=Inches(6.5), width=Inches(1), height=Inches(0.6), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], format_number_kmb(total_content_y), alignment=PP_ALIGN.CENTER, font_name= 'Poppins', font_size=14, font_bold=True,left=Inches(12), top=Inches(6.2), width=Inches(1), height=Inches(0.6), font_color=RGBColor(146, 39, 143))

	st.write("Slide 14 of 17")
	
#-----------PAGE 15---------------
	page_no = page_no + 1
	format_title(ppt.slides[page_no], "appendix - BRAND RANK", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=18, font_bold=True,left=Inches(0.5), top=Inches(0.3), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

	format_title(ppt.slides[page_no], "SOV", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_bold=True,left=Inches(0.5), top=Inches(0.8), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], "SOE", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_bold=True,left=Inches(4.5), top=Inches(0.8), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	format_title(ppt.slides[page_no], "SOC", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=14, font_bold=True,left=Inches(8.5), top=Inches(0.8), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))
	
# Add table	
	table_default(ppt.slides[page_no], rank_view[rank_view['Rank'] <= 20], Inches(0.5), Inches(1.2), Inches(4), Inches(7),
		      [Inches(0.9)]*4, Inches(0.3), header=True, upper=True, fontsize=9, alignment=PP_ALIGN.LEFT)	
	table_default(ppt.slides[page_no], rank_eng[rank_eng['Rank'] <= 20], Inches(4.5), Inches(1.2), Inches(4), Inches(7),
		      [Inches(0.9)]*4, Inches(0.3), header=True, upper=True, fontsize=9, alignment=PP_ALIGN.LEFT)
	table_default(ppt.slides[page_no], rank_content[rank_content['Rank'] <= 20], Inches(8.5), Inches(1.2), Inches(4), Inches(7),
		      [Inches(0.9)]*4, Inches(0.3), header=True, upper=True, fontsize=9, alignment=PP_ALIGN.LEFT)	
	
	st.write("Slide 15 of 17")
	
#------------PAGE 16--------------
	page_no = page_no + 1
	if 'Select All' in category_selection:
		category_title = "All"
	else: category_title = "| ".join(category_selection)
	
	format_title(ppt.slides[page_no], category_title.upper()+" CATEGORY KOL MIX", alignment=PP_ALIGN.LEFT, font_name= 'Neue Haas Grotesk Text Pro', font_size=28, font_bold=True,left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3), font_color=RGBColor(0, 0, 0))

# Filter the dataframe

	df_16 = df2[(df2['category'].isin(category)) & (df2['years'] == year) & (df2['month'].isin(month_num))]
	df_16_ = pd.pivot_table(df_16[['tier', 'views', 'engagements']], index = 'tier', values=['views', 'engagements'], aggfunc = 'sum', fill_value = 0)
	df_16_ = df_16_.sort_values(by=['tier'], ascending=False)
	df_16_ = df_16_[['views', 'engagements']]

# Add combo stacked bar chart
	combo2_chart(ppt.slides[page_no], df_16_, Inches(1), Inches(1.7), Inches(11), Inches(5), chart_title=True, title= f"{category_title} Category",
            fontsize=Pt(10), fontsize_title=Pt(12), smooth=True, data_show=True)
	
	st.write("Slide 16 of 17")
	
#-----------PAGE 17---------------
	page_no += 1
	format_title(ppt.slides[page_no], "Best Performing Content", alignment=PP_ALIGN.LEFT,
		     font_name='Montserrat', font_size=24, font_bold=True, 
		     left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(0.3),
		     font_color=RGBColor(0, 0, 0))
	
	format_title(ppt.slides[page_no], "Trending Content", alignment=PP_ALIGN.CENTER,
		     font_name='Montserrat', font_size=18, font_bold=True,
		     left=Inches(0.5), top=Inches(1.5), width=Inches(12.3), height=Inches(0.3),
		     font_color=RGBColor(0, 0, 0))
	
	format_title(ppt.slides[page_no], "(based on Engagement)", alignment=PP_ALIGN.CENTER,
		     font_name='Montserrat', font_size=12, font_bold=False, font_italic = True
		     left=Inches(0.5), top=Inches(1.25), width=Inches(12.3), height=Inches(0.3),
		     font_color=RGBColor(0, 0, 0))
	
	# Prepare data TOP 2 by er_content
	df_17 = df_16[['division', 'campaign', 'kol_name', 'link_post', 'views', 'engagements', 'er_content', 'followers']]
	df_17 = df_17.sort_values('er_content', ascending=False).head(2)
	
	# Prepare data TOP 2 by views
	df_17_ = df_16[(df_16['advocacy']== 'Boosted')]
	df_17_ = df_17_[['division', 'campaign', 'kol_name', 'link_post', 'views', 'engagements', 'er_content', 'followers']]
	df_17_ = df_17_.sort_values('views', ascending=False).head(2)
	
	df_17_join = pd.concat([df_17, df_17_], ignore_index=True)

	df_17_join['vr'] = ". . ."
	df_17_join['eng_score'] = np.where(df_17_join['followers'] != 0, df_17_join['engagements'] / df_17_join['followers'], 0)
	df_17_join['eng_score'] = np.ceil(df_17_join['eng_score'] * 100) / 100

	df_17_transpose = df_17_join[['kol_name', 'views', 'engagements', 'vr', 'er_content', 'eng_score', 'link_post']].transpose()
	df_17_transpose.reset_index(inplace=True)
	
# Add bullets text (kol name, pt 12, bold=True)
	format_title(ppt.slides[page_no], df_17_transpose.iloc[0,1], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=12, font_bold=True, 
		     left=Inches(1), top=Inches(5), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[0,2], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=12, font_bold=True, 
		     left=Inches(3.5), top=Inches(5), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[0,3], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=12, font_bold=True, 
		     left=Inches(7), top=Inches(5), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[0,4], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=12, font_bold=True, 
		     left=Inches(9.5), top=Inches(5), width=Inches(2), height=Inches(1))
# Add labeled_bullets text (kol name, pt 10, bold=True)
	df_to_labeled_bullets(ppt.slides[page_no], df_17_transpose.iloc[1:-1, [0,1]], Inches(1), Inches(5.2), Inches(2), Inches(1), font_size=10, font_bold = True)
	df_to_labeled_bullets(ppt.slides[page_no], df_17_transpose.iloc[1:-1, [0,2]], Inches(3.5), Inches(5.2), Inches(2), Inches(1), font_size=10, font_bold = True)
	df_to_labeled_bullets(ppt.slides[page_no], df_17_transpose.iloc[1:-1, [0,3]], Inches(7), Inches(5.2), Inches(2), Inches(1), font_size=10, font_bold = True)
	df_to_labeled_bullets(ppt.slides[page_no], df_17_transpose.iloc[1:-1, [0,4]], Inches(9.5), Inches(5.2), Inches(2), Inches(1), font_size=10, font_bold = True)
# Add bullets text (kol name, pt 9, bold=False)
	format_title(ppt.slides[page_no], df_17_transpose.iloc[-1,1], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=9, font_bold=False, 
		     left=Inches(1), top=Inches(6.75), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[-1,2], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=9, font_bold=False, 
		     left=Inches(3.5), top=Inches(6.75), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[-1,3], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=9, font_bold=False, 
		     left=Inches(7), top=Inches(6.75), width=Inches(2), height=Inches(1))
	format_title(ppt.slides[page_no], df_17_transpose.iloc[-1,4], alignment=PP_ALIGN.CENTER, font_name='Arial', font_size=9, font_bold=False, 
		     left=Inches(9.5), top=Inches(6.75), width=Inches(2), height=Inches(1))
	
	st.write("Slide 17 of 17")

#-----------END OF SLIDE---------------
	page_no = page_no + 1
	for i in range(page_no,len(ppt.slides)):
		try:
			xml_slides = ppt.slides._sldIdLst
			slides = list(xml_slides)
			xml_slides.remove(slides[page_no])
		except:
			pass

####### ------------- SAVE PPT -----------------
	#st.info("Saving the PPT please wait...")
	file = (f'{category_title.upper()} MONTHLY REPORT - {str(month)} {year}')
	filename = (f'{category_title.upper()} MONTHLY REPORT - {str(month)} {year}.pptx')
	files = ppt.save(filename)

# ✅ This code runs **only after clicking the submit button**
	st.success(f"✅ Your report is ready : '{category_title} Monthly Report - {str(month)} {year}.pptx'")

	# Store filename in session state for later use
	st.session_state["report_filename"] = filename

# Send Email Button (Only appears after submitting)
if "report_filename" in st.session_state:
	filename = st.session_state["report_filename"]
else:
	st.error("can't send email 📧 : No file generated yet. Please generate the report first.")
	st.stop()  # Stop execution if filename is missing
		
if st.button("Send Report", type= "secondary"):
	wib = pytz.timezone("Asia/Jakarta")
	now = datetime.now(wib)
	formatted_date = now.strftime("%Y-%m-%d %H:%M")
	
    # Gmail SMTP Configuration
	SMTP_SERVER = "smtp.gmail.com"
	SMTP_PORT = 587
	EMAIL_USER = "dyah.dinasari.groupm@gmail.com"
	EMAIL_PASS = "koxp pzgm ixws ihek"
    # Email Details
	send_to = email_list #["dyah.dinasari@groupm.com"]		
	subject = "[Test] m-Slide Goat-L'Oreal"
	body = """Hi team,
        
		This is a test email sent via Python SMTP. Please don't reply to this email!
  
  	Sending you the m-slide progress up until today.
  	This report is generated automatically with the following selection:
  	
   	year		: {}
   	month		: {}
    	quarter		: {}
     	category	: {}
      			{}
      	division	: {}	
       			{}
       	brand comparison: {}
   
	Regards,
	Dyah Dinasari""".format(year, month, quarter_, category_selection, category, division_selection, division, brands)

	# ✅ Create Email Message
	msg = MIMEMultipart()
	msg["From"] = EMAIL_USER
	msg["To"] = ", ".join(send_to)
	msg["Subject"] = subject
	msg.attach(MIMEText(body, "plain"))
    # ✅ Attach the PowerPoint File
	if os.path.exists(filename):  # Check if the file exists
		with open(filename, "rb") as file:
			part = MIMEApplication(file.read(), Name=os.path.basename(filename))
			part["Content-Disposition"] = f'attachment; filename="{os.path.basename(filename)}"'
			msg.attach(part)
	else:
		st.write(f"⚠ Warning: File '{filename}' not found. Email will be sent without attachment.")

    # ✅ Send Email via Gmail SMTP
	try:
		smtp = smtplib.SMTP(SMTP_SERVER, SMTP_PORT)
		smtp.starttls()  # Secure the connection
		smtp.login(EMAIL_USER, EMAIL_PASS)  # Login with App Password
		smtp.sendmail(EMAIL_USER, send_to, msg.as_string())  # Send email
		smtp.quit()
		st.success(f"✅ Email sent successfully! on: {formatted_date} please expect email delay around 10 minutes")
	except Exception as e:
		st.error(f"❌ Error: {e}")
			
#else:
	#st.warning("⚠ Please fill in the details and click 'Submit'.")
  
